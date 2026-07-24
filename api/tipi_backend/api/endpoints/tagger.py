import codecs
import logging
import os
import pickle
import re
import subprocess
import tempfile
from os.path import splitext
from typing import Annotated

from fastapi import APIRouter, File, Form, HTTPException, Query, UploadFile
from fastapi.responses import JSONResponse
from pdfminer.high_level import extract_text as extract_pdf_text
from docx import Document
from pptx import Presentation

from legal_segmenter import segment as segment_legal

import tipi_tasks
from tipi_backend.api import cache
from tipi_backend.api.business import get_tags, get_kbs
from tipi_backend.api.request_models import KbQuery
from tipi_backend.settings import Config


log = logging.getLogger(__name__)

router = APIRouter(prefix="/tagger", tags=["tagger"])


def filter_tags(result, kb):
    tags = result["result"]["tags"]
    new_topics = []
    new_tags = []
    for tag in tags:
        if tag["knowledgebase"] in kb:
            new_tags.append(tag)
            new_topics.append(tag["topic"])
    new_topics = list(set(new_topics))
    result["result"]["topics"] = new_topics
    result["result"]["tags"] = new_tags
    return result


def remove_fields(result):
    tags = result["result"]["tags"]
    for tag in tags:
        del tag["public"]


def units_with_tags(content, tags, kb):
    """Segmenta el texto (etapa 2) y devuelve `(total_unidades, unidades_con_tags)`.

    Cada unidad con tags incluye su `text` (necesario para la codificación
    NormTrace de la etapa 3). Las unidades sin coincidencias se descartan.
    """
    units = segment_legal(content)
    hit = []
    for unit in units:
        unit_result = tipi_tasks.tagger.extract_tags_from_text(unit.text, tags)
        unit_result = filter_tags(unit_result, kb)
        remove_fields(unit_result)
        unit_tags = unit_result["result"]["tags"]
        if not unit_tags:
            continue
        hit.append(
            {
                "unit_id": unit.unit_id,
                "unit_type": unit.unit_type,
                "number": unit.number,
                "heading": (unit.heading or "")[:200],
                "text": unit.text,
                "parent_id": unit.parent_id,
                "topics": unit_result["result"]["topics"],
                "tags": unit_tags,
            }
        )
    return len(units), hit


def segment_and_tag(content, tags, kb):
    """Bloque `segmentation` (etapa 2) para la respuesta: conteos por unidad, sin
    incluir el texto completo de cada unidad (para no inflar la respuesta)."""
    total, hit = units_with_tags(content, tags, kb)
    units_out = [{k: v for k, v in u.items() if k != "text"} for u in hit]
    return {
        "mode": "legal",
        "units_total": total,
        "units_with_tags": len(units_out),
        "units": units_out,
    }


# Extensiones de imagen que aceptamos como "foto o escaneo" (se leen por OCR).
_IMAGE_EXTS = (".jpg", ".jpeg", ".png", ".tif", ".tiff", ".bmp", ".webp", ".gif")


def _word_count(text: str) -> int:
    """Palabras "de verdad" (>=2 letras, con acentos) en el texto. Sirve para
    distinguir texto legible de ruido (números de página, artefactos)."""
    return len(re.findall(r"[A-Za-zÁÉÍÓÚÜÑáéíóúüñ]{2,}", text or ""))


def _looks_thin(text: str) -> bool:
    """¿El texto extraído es esencialmente vacío? (PDF escaneado sin capa de
    texto, o con una capa de basura). Umbral en palabras, no en caracteres: un
    escaneo puede traer >20 chars de ruido y aun así no tener contenido real."""
    return _word_count(text) < 6


def _sniff_kind(head: bytes, filename: str, mimetype: str) -> str:
    """Tipo real del archivo por firma binaria (magic bytes) primero, luego por
    extensión y content-type. No confiar solo en content-type: los navegadores y
    móviles suben PDFs válidos como 'application/octet-stream' (o sin tipo)."""
    ext = splitext(filename or "")[1].lower()
    mimetype = (mimetype or "").lower()
    # 1) Firmas binarias (lo más confiable).
    if head[:4] == b"%PDF":
        return "pdf"
    if head[:3] == b"\xff\xd8\xff" or head[:8] == b"\x89PNG\r\n\x1a\n" \
            or head[:4] in (b"II*\x00", b"MM\x00*") or head[:2] == b"BM" \
            or head[:6] in (b"GIF87a", b"GIF89a") \
            or (head[:4] == b"RIFF" and head[8:12] == b"WEBP"):
        return "image"
    if head[:8] == b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1":  # OLE2 (.doc, .ppt viejos)
        return "doc"
    if head[:2] == b"PK":  # zip: docx/pptx/xlsx
        if ext == ".pptx" or "presentation" in mimetype:
            return "pptx"
        return "docx"  # por defecto, el ofimático más común
    # 2) Extensión / content-type como respaldo.
    if ext == ".pdf" or mimetype == "application/pdf":
        return "pdf"
    if ext in _IMAGE_EXTS or mimetype.startswith("image/"):
        return "image"
    if ext == ".docx" or "wordprocessingml" in mimetype:
        return "docx"
    if ext == ".pptx" or "presentationml" in mimetype:
        return "pptx"
    if ext == ".doc" or mimetype == "application/msword":
        return "doc"
    if ext == ".txt" or mimetype == "text/plain":
        return "txt"
    return "unknown"


def _decode_text(raw: bytes) -> str:
    """Decodifica un .txt tolerando codificaciones no UTF-8 (Windows/Latin)."""
    for enc in ("utf-8", "utf-8-sig", "cp1252", "latin-1"):
        try:
            return raw.decode(enc).strip()
        except UnicodeDecodeError:
            continue
    return raw.decode("utf-8", errors="replace").strip()


def _ocr_pdf(path: str) -> str:
    """OCR de un PDF escaneado (sin capa de texto), como respaldo de pdfminer.

    Los binarios `tesseract`/`poppler` van en la imagen del api; las libs Python
    (`pdf2image`, `pytesseract`) son opcionales: si faltan, devuelve "" y el flujo
    normal reporta el error amigable de siempre. Tope de páginas y dpi configurables
    para acotar tiempo/memoria en documentos grandes.
    """
    try:
        from pdf2image import convert_from_path
        import pytesseract
    except ImportError:
        logging.warning("OCR de PDF no disponible: instala pdf2image y pytesseract.")
        return ""
    max_pages = int(os.environ.get("TAGGER_OCR_MAX_PAGES", "20"))
    dpi = int(os.environ.get("TAGGER_OCR_DPI", "200"))
    try:
        images = convert_from_path(path, dpi=dpi, first_page=1, last_page=max_pages)
    except Exception as exc:
        logging.warning("No se pudo rasterizar el PDF para OCR: %s", exc)
        return ""
    partes = []
    for img in images:
        try:
            partes.append(pytesseract.image_to_string(img, lang="spa"))
        except Exception as exc:
            logging.warning("Fallo de OCR en una página: %s", exc)
    return "\n".join(partes).strip()


def _ocr_image(path: str) -> str:
    """OCR directo de una imagen (foto o captura de pantalla). Mismo respaldo
    opcional que el de PDF; degrada a "" si faltan las libs o falla la lectura."""
    try:
        import pytesseract
        from PIL import Image
    except ImportError:
        logging.warning("OCR de imagen no disponible: instala pytesseract y Pillow.")
        return ""
    try:
        with Image.open(path) as img:
            return pytesseract.image_to_string(img, lang="spa").strip()
    except Exception as exc:
        logging.warning("Fallo de OCR de imagen: %s", exc)
        return ""


def _extract_text_from_file(file: UploadFile) -> str:
    """Extrae texto de un archivo subido de forma robusta.

    - Detecta el tipo por firma binaria (no solo por content-type), así un PDF
      subido como 'octet-stream' se lee igual.
    - PDF escaneado (sin capa de texto legible) -> OCR de respaldo.
    - Foto o imagen (jpg/png/tiff/…) -> OCR directo.
    El objetivo: si el documento tiene texto legible por una persona, lo sacamos.
    """
    raw = file.file.read()
    kind = _sniff_kind(raw[:16], file.filename or "", file.content_type or "")
    suffix = splitext(file.filename or "")[1] or (".pdf" if kind == "pdf" else "")
    with tempfile.NamedTemporaryFile(prefix="tipiscanner_", suffix=suffix) as f:
        f.write(raw)
        f.seek(0)
        if kind == "txt":
            text = _decode_text(raw)
        elif kind == "pdf":
            text = (extract_pdf_text(f.name) or "").strip()
            # PDF escaneado o con capa de texto pobre -> OCR; nos quedamos con el
            # resultado que traiga más palabras legibles.
            if _looks_thin(text):
                ocr_text = _ocr_pdf(f.name)
                if _word_count(ocr_text) > _word_count(text):
                    text = ocr_text
        elif kind == "image":
            text = _ocr_image(f.name)
        elif kind == "docx":
            doc = Document(f)
            text = "\n".join([para.text for para in doc.paragraphs]).strip()
        elif kind == "doc":
            result = subprocess.run(
                ["antiword", f.name], stdout=subprocess.PIPE, stderr=subprocess.PIPE
            )
            if result.returncode != 0:
                raise Exception(
                    f"Error al leer el archivo .doc: {result.stderr.decode('utf-8')}"
                )
            text = result.stdout.decode("utf-8").strip()
        elif kind == "pptx":
            ppt = Presentation(f)
            text = "\n".join(
                [
                    shape.text
                    for slide in ppt.slides
                    for shape in slide.shapes
                    if hasattr(shape, "text")
                ]
            ).strip()
        else:
            raise HTTPException(
                status_code=400,
                detail="Formato no soportado. Sube un PDF, una imagen (foto o escaneo), o un archivo .txt, .docx, .doc o .pptx.",
            )
    if not text or not text.strip():
        # Mensaje específico para foto/escaneo: distingue "no pudimos LEER el
        # archivo" de "lo leímos pero no coincidió con ninguna etiqueta".
        if kind in ("pdf", "image"):
            raise HTTPException(
                status_code=422,
                detail="No pudimos extraer texto del archivo. Si es una foto o un PDF escaneado, prueba con una imagen más nítida o un PDF con texto seleccionable.",
            )
        raise HTTPException(
            status_code=400,
            detail="Error al obtener el texto del fichero proporcionado. Pruebe con otro fichero.",
        )
    return text.strip()


@router.post("/")
def extract(
    text: Annotated[str, Form()] = "",
    file: Annotated[UploadFile | str | None, File()] = None,
    knowledgebase: Annotated[str, Form()] = "",
    segment: Annotated[str, Form()] = "",
    deep: Annotated[bool, Form()] = False,
):
    """Etiqueta el texto y devuelve los temas (ODS) y etiquetas que coinciden.

    Con `segment=legal` añade un bloque `segmentation` con los conteos por unidad
    jurídica (artículo, fracción, inciso, transitorio) del texto.

    Con `deep=true` encola además la codificación estructural NormTrace (etapa 3,
    asíncrona) de las unidades con tags y devuelve `normtrace_task_id`; el bloque
    `structural` se recupera en `GET /tagger/deep/{id}` cuando termina.
    """
    try:
        # Blank knowledgebase = no filter (all public KBs), matching Flask's default.
        # The empty default also stops Swagger "Try it out" from sending `"string"`.
        kb = get_kbs({"knowledgebase": knowledgebase or None})

        cache_key = Config.CACHE_TAGS
        tags = cache.get(cache_key)
        if tags is None:
            tags = get_tags()
            cache.set(cache_key, tags, timeout=5 * 60)
        tags = codecs.encode(pickle.dumps(tags), "base64").decode()
        tipi_tasks.init()

        content = ""
        if text:
            content = text
        elif isinstance(file, UploadFile):
            content = _extract_text_from_file(file)

        text_length = len(content.split())

        if text_length >= Config.TAGGER_MAX_WORDS:
            task = tipi_tasks.tagger.extract_tags_from_text.apply_async((content, tags))
            eta_time = int((text_length / 1000) * 4)
            return {
                "status": "PROCESSING",
                "task_id": task.id,
                "estimated_time": eta_time,
            }

        result = tipi_tasks.tagger.extract_tags_from_text(content, tags)
        result = filter_tags(result, kb)
        remove_fields(result)

        # Etapa 2/3 opcionales: segmentación jurídica y codificación NormTrace.
        if (segment == "legal" or deep) and content:
            total, hit = units_with_tags(content, tags, kb)
            result["segmentation"] = {
                "mode": "legal",
                "units_total": total,
                "units_with_tags": len(hit),
                "units": [{k: v for k, v in u.items() if k != "text"} for u in hit],
            }
            # Etapa 3: encola la codificación estructural (cola normtrace).
            if deep and hit:
                task = tipi_tasks.normtrace.analyze_units.apply_async((hit,))
                result["normtrace_task_id"] = task.id

        return result
    except HTTPException:
        raise
    except Exception as e:
        log.error(e)
        raise HTTPException(status_code=500, detail="Internal server error")


@router.get("/deep/{id}")
def normtrace_result(id: str):
    """Devuelve el bloque `structural` (codificación NormTrace) de la tarea deep.

    Mientras la tarea no termina responde `{status: PENDING|STARTED}`; al terminar,
    `{status: SUCCESS, structural: {...}}` con las unidades codificadas (cada una
    con `confidence_level` y `review_status`).
    """
    try:
        tipi_tasks.init()
        return tipi_tasks.normtrace.check_status_task(id)
    except Exception as e:
        log.error(e)
        return JSONResponse(status_code=404, content={"Error": "No task found"})


@router.get("/result/{id}")
def tagger_result(id: str, query: Annotated[KbQuery, Query()]):
    """Devuelve el resultado de la tarea de etiquetado asíncrona."""
    try:
        tipi_tasks.init()
        result = tipi_tasks.tagger.check_status_task(id)

        kb = get_kbs(query.model_dump())
        if result["status"] == "SUCCESS":
            result = filter_tags(result, kb)
            remove_fields(result)
        return result
    except Exception as e:
        log.error(e)
        return JSONResponse(status_code=404, content={"Error": "No task found"})
